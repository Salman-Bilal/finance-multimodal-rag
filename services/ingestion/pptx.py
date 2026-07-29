from pptx import Presentation

def extract_pptx_chunks(file_path: str) -> list[str]:
    """Extract text from PowerPoint (.pptx) slides."""
    prs = Presentation(file_path)
    slide_chunks = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        if slide_text:
            combined_slide = f"Slide {slide_number}:\n" + "\n".join(slide_text)
            slide_chunks.append(combined_slide)

    return slide_chunks if slide_chunks else ["PowerPoint presentation contains no text content."]